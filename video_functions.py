import os
import imageio
import convertapi
import requests
from pptx import Presentation
from moviepy.editor import *
import firebase_admin
from firebase_admin import credentials, storage
from moviepy.editor import *
import uuid
from dotenv import load_dotenv

load_dotenv()

DID_API_KEY = os.getenv("DID_API_KEY")


def uploadToFirebase(filepath):
    bucket = storage.bucket()
    blob = bucket.blob(uuid.uuid4().hex)
    blob.upload_from_filename(filepath)
    blob.make_public()
    return blob.public_url    

# Your existing functions with added docstrings
def concatenate_videos(video_bytes_list, output_path='concatenated.mp4'):
    """
    Concatenates a list of video bytes into a single video file.

    Args:
        video_bytes_list (list): A list of video bytes.
        output_path (str, optional): The output path for the concatenated video file. 
            Defaults to 'concatenated.mp4'.

    Returns:
        str: The output path of the concatenated video file.
    """
    filenames = [f"toconcat{i}.mp4" for i in range(len(video_bytes_list))]
    for i,filename in enumerate(filenames):
        with open(filename, "wb") as file:
            file.write(requests.get(video_bytes_list[i]).content)

    clips = [VideoFileClip(filename) for filename in filenames]
    final = concatenate_videoclips(clips)
    final.write_videofile(output_path)
    return output_path

def read_images_as_bytes(image_paths):
    """
    Read a list of image files and return their contents as bytes.

    Args:
        image_paths (list): A list of file paths to the image files.

    Returns:
        list: A list of image contents as bytes.

    Raises:
        FileNotFoundError: If an image file is not found.
        IOError: If there is an error reading an image file.
    """
    image_bytes_list = []

    for image_path in image_paths:
        try:
            # Read the image file as bytes
            with open(image_path, 'rb') as image_file:
                image_bytes = image_file.read()
                image_bytes_list.append(image_bytes)

        except FileNotFoundError as e:
            print(f"Error reading image {image_path}: File not found.")
            raise e

        except IOError as e:
            print(f"Error reading image {image_path}: {e}")
            raise e

    return image_bytes_list

def generate_video(img_paths, durations_seconds, output_path='output.mp4'):
    """
    Generate a video by concatenating images with specified durations.

    Parameters:
    - img_paths (list): List of paths to image files.
    - durations_seconds (list): List of durations (in seconds) for each image.
    - output_path (str): Output path for the generated video.

    Returns:
    - str: Path to the generated video.
    """
    images = []

    for img_path, duration in zip(img_paths, durations_seconds):
        # Read image
        img = imageio.imread(img_path)

        # Duplicate the image to display for the specified duration
        frames = [img] * int(duration * 10)  # 10 frames per second

        # Add frames to the video
        images.extend(frames)

    # Write the video
    output_path = os.path.abspath(output_path)
    imageio.mimsave(output_path, images, 'mp4', fps=10)  # Set the desired frames per second

    return output_path

def read_pptx(file_path):
    """
    Read a PowerPoint presentation and extract notes and images.

    Parameters:
    - file_path (str): Path to the PowerPoint presentation file.

    Returns:
    - Tuple: Tuple containing:
        - List: Commenter notes for each slide.
        - List: Paths to images extracted from the presentation.
        - List: Texts extracted from each shape in the presentation.
    """
    # Load the PowerPoint presentation
    presentation = Presentation(file_path)

    commenter_notes = []  # List to store commenter notes
    slide_txts = []

    # Generate images
    convertapi.api_secret = 'yzCX1crd9TDqqdRZ'
    save_path = 'imgs'
    imgs_paths = convertapi.convert('jpg', {
        'File': file_path
    }, from_format='pptx').save_files(save_path)

    # Iterate through each slide
    for i, slide in enumerate(presentation.slides):
        slide_notes = []  # List to store commenter notes for the current slide

        # Extract text from each shape in the slide
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text
                slide_txts.append(text)

        # Extract comments from the slide
        comment = slide.notes_slide.notes_text_frame
        text = comment.text
        slide_notes.append(text)

        commenter_notes.append(slide_notes)

        # Extract images from shapes
        # image = slide_to_image(slide)
        # slide_images.append(image)

    return commenter_notes, imgs_paths, slide_txts

def get_presentation_video(pptx_file_path, durations_seconds):
    """
    Generate a video from a PowerPoint presentation.

    Parameters:
    - pptx_file_path (str): Path to the PowerPoint presentation file.
    - durations_seconds (list): List of durations (in seconds) for each slide.

    Returns:
    - str: Path to the generated video.
    """
    notes, img_paths, slide_txts = read_pptx(pptx_file_path)
    video_path = generate_video(img_paths, durations_seconds)
    return video_path

import requests

def avatarVideoRequest(avatar_img_url, text, speech_motion, gender):
    """
    Sends a request to generate an avatar video with the given parameters.

    Args:
        avatar_img_url (str): The URL of the avatar image.
        text (str): The text to be spoken in the video.
        speech_motion (str): The style of speech motion for the avatar. Can be "friendly" or "cheerful".
        gender (str): The gender of the avatar. Can be "male" or "female".

    Returns:
        str: The ID of the generated video.

    Raises:
        requests.exceptions.RequestException: If the request fails.

    """
    url = "https://api.d-id.com/talks"
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Basic {DID_API_KEY}",
    }
    data = {
        "script": {
            "type": "text",
            "input": text,
            "provider": {
                "type": "microsoft",
                "voice_id": "en-US-TonyNeural" if gender == "male" else "en-US-JennyNeural",
                "voice_config": {
                    "style": speech_motion
                }
            }
        },
        "source_url": avatar_img_url,
        "config": {
            "driver_expressions": {
                "expressions": [
                    {
                        "start_frame": 0,
                        "expression": "happy" if speech_motion in ["friendly", "cheerful"] else "neutral",
                        "intensity": 1
                    }
                ]
            }
        }
    }

    response = requests.post(url, json=data, headers=headers)

    # Check the response
    if response.status_code == 200:
        print("Request was successful.")
        print(response.json())
    else:
        print(f"Request failed with status code {response.status_code}:")
        print(response.text)
    
    id = response.json()['id']
    return id

def getAvatarVideo(id):
    url = f"https://api.d-id.com/talks/{id}"
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Basic {DID_API_KEY}",
    }

    response = requests.get(url, headers=headers)

    # Check the response
    if response.status_code == 200:
        print("Request was successful.")
        print(response.json())
    else:
        print(f"Request failed with status code {response.status_code}:")
        print(response.text)
    videourl = response.json()['result_url']
    
    duration_seconds = response.json()['duration']
    return videourl,duration_seconds


def convertUrl(request):
    temp  =request.avatar_img_url.split("uploads%2F")[1].split("?")[0]
    avatar_img_url = f"https://storage.googleapis.com/vectordatabase-95677.appspot.com/uploads/{temp}"
    return avatar_img_url